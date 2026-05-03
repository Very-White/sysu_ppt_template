"""
将 HTML PPT 模板导出为 PPTX 文件。
用法：
    conda activate sysu-ppt-export
    python export_to_pptx.py
    python export_to_pptx.py -i my_slides.html -o presentation.pptx
"""

import argparse
import os
import sys
import tempfile
import shutil
import threading
import http.server
import socketserver
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu


SLIDE_WIDTH = 1280
SLIDE_HEIGHT = 720


def _find_free_port():
    with socketserver.TCPServer(("", 0), None) as s:
        return s.server_address[1]


def _start_server(directory: str, port: int):
    handler = http.server.SimpleHTTPRequestHandler
    os.chdir(directory)
    server = socketserver.TCPServer(("", port), handler)
    server.serve_forever()


def capture_slides(html_path: str, output_dir: str):
    html_path = os.path.abspath(html_path)
    if not os.path.isfile(html_path):
        print(f"错误：找不到文件 {html_path}")
        sys.exit(1)

    html_dir = os.path.dirname(html_path)
    html_file = os.path.basename(html_path)

    port = _find_free_port()
    server_thread = threading.Thread(
        target=_start_server, args=(html_dir, port), daemon=True
    )
    server_thread.start()

    file_url = f"http://localhost:{port}/{html_file}"
    print(f"本地 HTTP 服务器已启动: {file_url}")

    screenshots = []

    with sync_playwright() as p:
        browser = None
        for ch in ("msedge", "chrome"):
            try:
                browser = p.chromium.launch(channel=ch)
                print(f"使用系统 {ch}")
                break
            except Exception:
                continue
        if browser is None:
            browser = p.chromium.launch()
            print("使用 Playwright 内置 Chromium")
        page = browser.new_page(
            viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT},
            device_scale_factor=2,
        )
        page.goto(file_url)
        page.wait_for_load_state("networkidle")
        page.wait_for_timeout(2000)

        slide_count = page.evaluate(
            "() => document.querySelectorAll('#slide-deck .slide').length"
        )
        print(f"检测到 {slide_count} 张幻灯片")

        page.keyboard.press("Home")
        page.wait_for_timeout(500)

        page.evaluate(
            "() => {"
            "  const d = document.getElementById('slide-deck');"
            "  d.style.transform = 'none';"
            "  d.style.position = 'relative';"
            "  d.style.left = '0';"
            "  d.style.top = '0';"
            "  d.style.width = '1280px';"
            "  d.style.height = '720px';"
            "}"
        )

        for i in range(slide_count):
            if i > 0:
                page.keyboard.press("ArrowRight")
                page.wait_for_timeout(500)

            deck = page.locator("#slide-deck")

            filename = os.path.join(output_dir, f"slide_{i:03d}.png")
            deck.screenshot(path=filename)
            screenshots.append(filename)
            print(f"  已截图第 {i + 1}/{slide_count} 页")

        browser.close()

    return screenshots


def create_pptx(screenshots: list, output_path: str):
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for img_path in screenshots:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(
            img_path, Emu(0), Emu(0), slide_width, slide_height
        )

    prs.save(output_path)
    print(f"PPTX 已保存至: {output_path}")


def main():
    parser = argparse.ArgumentParser(description="将 HTML PPT 导出为 PPTX")
    parser.add_argument(
        "-i", "--input", default="index.html", help="输入 HTML 文件路径"
    )
    parser.add_argument(
        "-o", "--output", default="output.pptx", help="输出 PPTX 文件路径"
    )
    args = parser.parse_args()

    script_dir = os.path.dirname(os.path.abspath(__file__))
    input_path = os.path.join(script_dir, args.input) if not os.path.isabs(args.input) else args.input
    output_path = os.path.join(script_dir, args.output) if not os.path.isabs(args.output) else args.output

    tmp_dir = tempfile.mkdtemp(prefix="sysu_ppt_")
    try:
        print("正在截图...")
        screenshots = capture_slides(input_path, tmp_dir)
        print(f"共截图 {len(screenshots)} 张，正在生成 PPTX...")
        create_pptx(screenshots, output_path)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


if __name__ == "__main__":
    main()
